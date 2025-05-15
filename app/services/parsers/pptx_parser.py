# app/services/parsers/ppt_parser.py

from pptx import Presentation

def parse_ppt(file_path: str) -> str:
    """
    Extracts text and table contents from a PPTX file.
    """
    extracted_text = ""

    presentation = Presentation(file_path)

    for i, slide in enumerate(presentation.slides):
        slide_text = f"--- Slide {i + 1} ---\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"

            # Handle tables inside slides
            if shape.shape_type == 19:  # shape_type 19 = TABLE
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    slide_text += row_text + "\n"
        
        extracted_text += slide_text + "\n"

    return extracted_text.strip()
